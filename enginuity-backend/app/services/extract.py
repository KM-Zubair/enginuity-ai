# app/services/extract.py
from pathlib import Path
from typing import Tuple
import fitz  
from pptx import Presentation

def from_pdf(path: Path) -> str:
    doc = fitz.open(str(path))
    texts = []
    for page in doc:
        texts.append(page.get_text()) # type: ignore
    return "\n".join(texts).strip()

def from_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                buf.append(shape.text)
        if buf:
            texts.append("\n".join(buf))
    return "\n\n".join(texts).strip()

def extract_text(path: Path) -> Tuple[str, str]:
    """Returns (title, full_text)."""
    p = Path(path)
    name = p.stem
    if p.suffix.lower() == ".pdf":
        return name, from_pdf(p)
    elif p.suffix.lower() == ".pptx":
        return name, from_pptx(p)
    else:
        # audio pipeline can put its transcript.txt here later
        return name, p.read_text(encoding="utf-8", errors="ignore")
